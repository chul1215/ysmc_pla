import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "대전 지역 성형외과 경쟁 벤치마킹 분석"
slide.placeholders[1].text = "유성선병원 성형외과 포지셔닝 전략\n2026.03."

# Slide 2: 시장 환경 및 경쟁 지형
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "대전 성형외과 시장 및 경쟁 지형"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "둔산동 중심의 밀집 상권"
p = tf.add_paragraph(); p.text = "대부분의 의원이 둔산/탄방동에 밀집하여 가격 및 홍보 경쟁 심화"; p.level = 1
p = tf.add_paragraph(); p.text = "주요 경쟁자 현황"; p.level = 0
p = tf.add_paragraph(); p.text = "닥터스미, 페이스 등 규모 우위 의원 (미용 직접 경쟁)"; p.level = 1
p = tf.add_paragraph(); p.text = "오체안 대전점 (치료/재건/외상 상권 전문화로 경쟁)"; p.level = 1
p = tf.add_paragraph(); p.text = "을지대병원 (미용+재건, 동일 포지셔닝 위협)"; p.level = 1
p = tf.add_paragraph(); p.text = "유성구 내 기회 창출"; p.level = 0
p = tf.add_paragraph(); p.text = "지역 내 1개 의원만 존재하여 수요 공백 및 상권 선점 가능성"; p.level = 1

# Slide 3: SWOT
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "유성선병원 성형외과 SWOT 분석"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "강점 (Strengths)"
p = tf.add_paragraph(); p.text = "종합병원 인프라 완비 (마취/응급), 미용~재건 풀커버 1:1 맞춤의료진"; p.level = 1
p = tf.add_paragraph(); p.text = "약점 (Weaknesses)"; p.level = 0
p = tf.add_paragraph(); p.text = "초기 신설 부서 인지도 부족 및 둔산동 거리 접근성 이슈"; p.level = 1
p = tf.add_paragraph(); p.text = "기회 (Opportunities)"; p.level = 0
p = tf.add_paragraph(); p.text = "안전 중심의 트렌드, 유성구 진료 수요, 기존 선병원 타과 연계"; p.level = 1
p = tf.add_paragraph(); p.text = "위협 (Threats)"; p.level = 0
p = tf.add_paragraph(); p.text = "둔산동 대형 의원의 유출 마케팅 강화, 대학병원(충남대/을지대) 포지션 경쟁"; p.level = 1

# Slide 4: 포지셔닝 & 핵심 메시지
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "차별화 포지셔닝 및 전략"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "메인 슬로건: '대학병원의 안전함에 섬세함을 더하다'"
p = tf.add_paragraph(); p.text = "초-안전 프리미엄 전략"; p.level = 0
p = tf.add_paragraph(); p.text = "단순 가격 경쟁 회피, '마취과 전문의 상주 및 중환자실 인프라' 적극 강조"; p.level = 1
p = tf.add_paragraph(); p.text = "공장식 성형 탈피"; p.level = 0
p = tf.add_paragraph(); p.text = "가톨릭중앙의료원 출신 전담 전문의 직접 진료-수술-경과관리"; p.level = 1
p = tf.add_paragraph(); p.text = "원스톱 토탈 케어 솔루션 구축"; p.level = 0
p = tf.add_paragraph(); p.text = "미용부터 재건까지 하나의 과목에서 해결"; p.level = 1

# Slide 5: 마케팅 커뮤니케이션 로드맵
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "단계별 실행 로드맵 (Action Plan)"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "1단계: 노출 및 인지도 구축 (1~3개월)"
p = tf.add_paragraph(); p.text = "네이버 플레이스 최적화, 공식 블로그 초기 정보 포스팅"; p.level = 1
p = tf.add_paragraph(); p.text = "2단계: 예약 유입 확장 (4~6개월)"; p.level = 0
p = tf.add_paragraph(); p.text = "비경쟁 틈새/안전 키워드 타겟 광고, 카카오톡 1:1 예약 채널 도입"; p.level = 1
p = tf.add_paragraph(); p.text = "3단계: 브랜딩 자산 고도화 (7~12개월)"; p.level = 0
p = tf.add_paragraph(); p.text = "유튜브/인스타그램 확장 전개, 리뷰 및 후기 컨텐츠 확보"; p.level = 1
p = tf.add_paragraph(); p.text = "4단계: 시장 선도 입지 다지기 (12개월 이후)"; p.level = 0
p = tf.add_paragraph(); p.text = "외상/재건 특화 지역 내 탑티어로 자리매김"; p.level = 1

prs.save('/Users/chul/Documents/WORK/ysmc_pla/briefing/competitor_analysis_presentation.pptx')
